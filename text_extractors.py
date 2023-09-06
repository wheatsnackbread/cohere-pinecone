from PyPDF2 import PdfReader  # Make sure to import the new class
from docx import Document
from pptx import Presentation


def read_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text + " "
                text += "\n"
    return text.strip()


def read_word(file_path):
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += (
            para.text.replace("\n", " ").strip() + " "
        )  # Replace newlines with space and append
    return text.strip()  # Remove trailing space at the end


def read_pdf(file_path):
    with open(file_path, "rb") as f:
        reader = PdfReader(f)  # Use the new PdfReader class
        text = ""
        for i in range(len(reader.pages)):
            page_text = reader.pages[i].extract_text()
            text += (
                page_text.replace("\n", " ").strip() + " "
            )  # Replace newlines with space and append
    return text.strip()  # Remove trailing space at the end


# fileName = "Sam'sFirstStockPitch.pptx"
# documentText = ""

# if fileName.endswith("pdf"):
#     documentText = read_pdf(fileName)
#     print(documentText)
# elif fileName.endswith("docx"):
#     documentText = read_word(fileName)
#     print(documentText)
# elif fileName.endswith('pptx'):
#     documentText = read_pptx(fileName)
#     print(documentText)
